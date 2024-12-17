import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from pdf2image import convert_from_bytes

# Streamlit 애플리케이션 제목
st.title("PDF to PPTX 변환기")

# PDF 파일 업로드
uploaded_file = st.file_uploader("PDF 파일을 업로드하세요", type=["pdf"])

# PDF -> PPTX 변환 함수
def pdf_to_pptx(pdf_bytes):
    # PDF 파일을 이미지로 변환
    images = convert_from_bytes(pdf_bytes)
    
    # PowerPoint 프레젠테이션 객체 생성
    presentation = Presentation()
    
    # 각 페이지를 슬라이드로 추가
    for img in images:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 빈 슬라이드 레이아웃
        left = top = Inches(0)
        pic = slide.shapes.add_picture(BytesIO(img.tobytes()), left, top, width=Inches(10), height=Inches(7.5))
    
    # PPTX를 메모리 바이트에 저장
    pptx_io = BytesIO()
    presentation.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# 파일 처리 및 다운로드 제공
if uploaded_file:
    st.info("파일 변환 중입니다. 잠시만 기다려주세요...")
    try:
        # PDF를 PPTX로 변환
        pptx_file = pdf_to_pptx(uploaded_file.read())
        
        # PPTX 다운로드 버튼 제공
        st.success("PDF 파일이 성공적으로 PPTX로 변환되었습니다!")
        st.download_button(
            label="PPTX 다운로드",
            data=pptx_file,
            file_name="converted_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.error(f"오류가 발생했습니다: {e}")
